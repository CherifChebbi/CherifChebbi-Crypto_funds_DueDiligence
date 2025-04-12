# chatbot_text_extraction.py

import os
import re
from pathlib import Path
import fitz  # PyMuPDF
import docx
import pytesseract
import pdfplumber
from bs4 import BeautifulSoup
from pdf2image import convert_from_path
from PIL import Image
import openpyxl  # Pour Excel
from pptx import Presentation  # Pour PowerPoint
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords
import nltk
import filetype
import csv
from unidecode import unidecode
import string
from spellchecker import SpellChecker

# Téléchargement des ressources nécessaires pour nltk
nltk.download('punkt')
nltk.download('stopwords')

# Créer les répertoires si pas encore faits
UPLOAD_DIR = "upload"
OUTPUT_DIR = "output"
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Fonction pour extraire le texte d'un fichier PDF
def extract_text_from_pdf(file_path):
    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""
        # Si l'extraction est insuffisante, essayer l'OCR
        if len(text.strip()) < 100:
            print("🧾 PDF probablement scanné – passage en OCR.")
            text = ocr_pdf(file_path)
    except Exception as e:
        print(f"❌ Erreur PDF classique : {e} – fallback OCR.")
        text = ocr_pdf(file_path)
    return text

# OCR pour PDF (scanner les images en texte)
def ocr_pdf(file_path):
    text = ""
    try:
        images = convert_from_path(file_path)
        for image in images:
            text += pytesseract.image_to_string(image)
    except Exception as e:
        print(f"❌ Erreur OCR PDF : {e}")
    return text

# Extraction de texte à partir de fichiers DOCX
def extract_text_from_docx(file_path):
    text = ""
    try:
        doc = docx.Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        print(f"❌ Erreur DOCX : {e}")
    return text

# Extraction de texte à partir de fichiers Excel (XLSX)
def extract_text_from_excel(file_path):
    text = ""
    try:
        wb = openpyxl.load_workbook(file_path, read_only=True)
        for sheet in wb.sheetnames:
            sheet_data = wb[sheet]
            for row in sheet_data.iter_rows(values_only=True):
                text += " ".join([str(cell) for cell in row]) + "\n"
    except Exception as e:
        print(f"❌ Erreur Excel : {e}")
    return text

# Extraction de texte à partir de fichiers PowerPoint (PPTX)
def extract_text_from_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"❌ Erreur PPTX : {e}")
    return text

# Extraction de texte à partir de fichiers texte (TXT)
def extract_text_from_txt(file_path):
    try:
        with open(file_path, 'r', encoding="utf-8") as f:
            return f.read()
    except Exception as e:
        print(f"❌ Erreur TXT : {e}")
        return ""

# Extraction de texte à partir de fichiers HTML
def extract_text_from_html(file_path):
    try:
        with open(file_path, 'r', encoding="utf-8") as f:
            soup = BeautifulSoup(f.read(), "html.parser")
            return soup.get_text(separator="\n")
    except Exception as e:
        print(f"❌ Erreur HTML : {e}")
        return ""

# Extraction de texte à partir de fichiers CSV
def extract_text_from_csv(file_path):
    text = ""
    try:
        with open(file_path, 'r', encoding="utf-8") as f:
            reader = csv.reader(f)
            for row in reader:
                text += " ".join(row) + "\n"
    except Exception as e:
        print(f"❌ Erreur CSV : {e}")
    return text

# OCR pour les images (PNG, JPG, etc.)
def ocr_image(file_path):
    text = ""
    try:
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
    except Exception as e:
        print(f"❌ Erreur OCR image : {e}")
    return text


# 📌 Fonction pour nettoyer le texte extrait
def clean_text(text):
    print("🧼 Nettoyage du texte...")
    text = re.sub(r"\n{2,}", "\n", text)  # Supprimer les lignes vides multiples
    text = re.sub(r"\s{2,}", " ", text)   # Supprimer les espaces multiples
    text = text.replace("\x0c", "")       # Supprimer les caractères spéciaux
    return text.strip()


# Correction orthographique du texte
def correct_spelling(text):
    spell = SpellChecker()
    words = text.split()
    corrected_text = ' '.join([spell.correction(word) for word in words])
    return corrected_text

# Détection du type de fichier (extension)
def detect_file_type(file_path):
    kind = filetype.guess(file_path)
    ext = kind.extension if kind else Path(file_path).suffix[1:].lower()
    return ext

# Fonction principale pour extraire le texte selon le type de fichier
def extract_text(file_path):
    ext = detect_file_type(file_path)
    print(f"📄 Traitement du fichier : {file_path} (type détecté : {ext})")
    
    if ext == "pdf":
        text = extract_text_from_pdf(file_path)
    elif ext == "docx":
        text = extract_text_from_docx(file_path)
    elif ext == "xlsx":
        text = extract_text_from_excel(file_path)
    elif ext == "pptx":
        text = extract_text_from_pptx(file_path)
    elif ext == "txt":
        text = extract_text_from_txt(file_path)
    elif ext in ["html", "htm"]:
        text = extract_text_from_html(file_path)
    elif ext == "csv":
        text = extract_text_from_csv(file_path)
    elif ext in ["png", "jpg", "jpeg"]:
        text = ocr_image(file_path)
    else:
        print("⚠️ Format non supporté.")
        return None

    return clean_text(text)

# Sauvegarde du texte extrait dans un fichier
def save_extracted_text(file_path, text):
    file_stem = Path(file_path).stem
    file_stem = re.sub(r'[<>:"/\\|?*]', '_', file_stem)  # Remplace les caractères invalides
    file_stem = file_stem.replace(" ", "_")  # Remplace les espaces par des underscores
    out_dir = Path(OUTPUT_DIR) / file_stem
    out_dir.mkdir(parents=True, exist_ok=True)
    
    out_path = out_dir / "extracted_text.txt"
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(text)
    print(f"✅ Texte extrait sauvegardé dans : {out_path}")
    
    # Appel de la fonction d'évaluation
    evaluate_text_quality(text, file_stem)

# Traitement de tous les fichiers dans le répertoire d'upload
def process_all_files():
    for filename in os.listdir(UPLOAD_DIR):
        file_path = os.path.join(UPLOAD_DIR, filename)
        if os.path.isfile(file_path):
            text = extract_text(file_path)
            if text and len(text.strip()) > 50:
                save_extracted_text(file_path, text)
            else:
                print(f"❌ Extraction échouée ou vide : {filename}")

# Évaluation de la qualité du texte extrait
def evaluate_text_quality(text, filename=""):
    print(f"\n📊 Évaluation du texte extrait : {filename}")

    # Tokenisation des phrases et des mots
    sentences = sent_tokenize(text)
    words = word_tokenize(text)
    stop_words = set(stopwords.words('english'))

    # Nombre de phrases, mots, et mots non-communs (sans stopwords)
    num_sentences = len(sentences)
    num_words = len(words)
    num_non_stopwords = len([w for w in words if w.lower() not in stop_words and w.isalpha()])
    num_empty_lines = text.count('\n\n')
    
    # Fréquence des mots clés financiers (à adapter selon votre domaine)
    financial_keywords = ['investment', 'fund', 'return', 'risk', 'portfolio', 'strategy']
    words_lower = [word.lower() for word in words if word.isalpha()]
    keyword_counts = {kw: words_lower.count(kw) for kw in financial_keywords}

    # Analyser la ponctuation et les symboles
    punctuation_count = sum(1 for char in text if char in string.punctuation)

    print(f"📝 Nombre de phrases : {num_sentences}")
    print(f"🔤 Nombre total de mots : {num_words}")
    print(f"💡 Nombre de mots utiles (non stopwords) : {num_non_stopwords}")
    print(f"📉 Pourcentage de lignes vides : {round((num_empty_lines / (text.count('\n') + 1)) * 100, 2)}%")
    print(f"📚 Densité de mots utiles : {round((num_non_stopwords / num_words) * 100, 2)}%")
    print(f"📐 Longueur moyenne par phrase : {round(num_words / num_sentences, 2)} mots")
    print(f"🔠 Nombre de ponctuations : {punctuation_count}")
    
    # Fréquence des mots-clés
    print("🔑 Fréquence des mots-clés financiers :")
    for keyword, count in keyword_counts.items():
        print(f"{keyword}: {count}")

    # Détection de sections structurées
    common_sections = ['executive summary', 'contact', 'fund overview', 'risk factors', 'disclaimer']
    structure_hits = [sec for sec in common_sections if sec in text.lower()]
    if structure_hits:
        print(f"✅ Sections détectées : {', '.join(structure_hits)}")
    else:
        print("⚠️ Aucune section structurée détectée.")

    # Vérification de la cohérence des mots-clés
    missing_keywords = [kw for kw in financial_keywords if kw not in keyword_counts or keyword_counts[kw] == 0]
    if missing_keywords:
        print(f"⚠️ Mots-clés manquants : {', '.join(missing_keywords)}")
    else:
        print("✅ Tous les mots-clés sont présents.")

if __name__ == "__main__":
    print("🚀 Démarrage de l'extraction de texte multi-formats...")
    process_all_files()
    print("🎯 Extraction terminée.")
